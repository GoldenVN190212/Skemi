# Train/extract_universal.py
import os
import json
import zipfile
import sqlite3
import logging
import tempfile
import re

# Documents
import docx
import pdfplumber
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
import yaml

# Images
from PIL import Image
import pytesseract

# Audio / Video
import speech_recognition as sr
from pydub import AudioSegment
from moviepy import VideoFileClip

# Optional: Magic
try:
    import magic
    USE_MAGIC = True
except ImportError:
    USE_MAGIC = False

logging.basicConfig(level=logging.INFO)

# -------------------------
# Helper: Extract binary text as fallback
# -------------------------
def extract_from_binary(path):
    try:
        raw = open(path, "rb").read()
        ascii_text = re.findall(rb"[ -~]{5,}", raw)
        unicode_text = re.findall(rb"(?:[\x20-\x7E][\x00]){5,}", raw)
        ascii_text = b"\n".join(ascii_text).decode("utf-8", errors="ignore")
        unicode_text = b"\n".join(unicode_text).decode("utf-16", errors="ignore")
        return ascii_text + "\n" + unicode_text
    except:
        return ""

# -------------------------
# Extractors: Documents
# -------------------------
def extract_pdf(path):
    try:
        text = ""
        with pdfplumber.open(path) as pdf:
            for p in pdf.pages:
                t = p.extract_text()
                if t:
                    text += t + "\n"
        return text
    except Exception as e:
        logging.error(f"PDF Error: {e}")
        return ""

def extract_docx(path):
    try:
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logging.error(f"DOCX Error: {e}")
        return ""

def extract_pptx(path):
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        logging.error(f"PPTX Error: {e}")
        return ""

def extract_excel(path):
    try:
        df_dict = pd.read_excel(path, sheet_name=None)
        text = ""
        for sheet_name, df in df_dict.items():
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.to_string()
        return text
    except Exception as e:
        logging.error(f"Excel Error: {e}")
        return ""

def extract_csv(path):
    try:
        df = pd.read_csv(path)
        return df.to_string()
    except Exception as e:
        logging.error(f"CSV Error: {e}")
        return ""

# -------------------------
# Extractors: Structured / Code
# -------------------------
def extract_json(path):
    try:
        data = json.load(open(path, "r", encoding="utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except:
        return ""

def extract_yaml(path):
    try:
        data = yaml.safe_load(open(path, "r", encoding="utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except:
        return ""

def extract_html(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text("\n")
    except:
        return ""

def extract_txt(path):
    try:
        return open(path, "r", encoding="utf-8", errors="ignore").read()
    except:
        return ""

def extract_sqlite(path):
    try:
        conn = sqlite3.connect(path)
        cursor = conn.cursor()
        output = []
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
        tables = cursor.fetchall()
        for (table,) in tables:
            cursor.execute(f"SELECT * FROM {table}")
            rows = cursor.fetchall()
            output.append(f"===== TABLE {table} =====")
            for row in rows:
                output.append(str(row))
        conn.close()
        return "\n".join(output)
    except:
        return ""

def extract_zip(path):
    text = ""
    try:
        with zipfile.ZipFile(path) as z:
            for file in z.namelist():
                if file.startswith("__") or file.endswith("/"):
                    continue
                try:
                    with z.open(file) as sub:
                        content = sub.read().decode("utf-8", errors="ignore")
                        text += f"\n\n=== {file} ===\n{content}"
                except:
                    continue
    except:
        pass
    return text

# -------------------------
# Extractors: Images
# -------------------------
def extract_image(path):
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img, lang="vie+eng").strip()
    except:
        return ""

# -------------------------
# Extractors: Audio / Video
# -------------------------
def extract_audio(path):
    recognizer = sr.Recognizer()
    wav_path = None
    try:
        audio = AudioSegment.from_file(path)
        if len(audio) > 120000:
            audio = audio[:120000]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            wav_path = tmp.name
            audio.export(wav_path, format="wav")
        with sr.AudioFile(wav_path) as source:
            audio_data = recognizer.record(source)
            text = recognizer.recognize_google(audio_data, language="vi-VN")
            return f"[Audio]: {text}"
    except:
        return "[Audio]: Không thể nhận diện"
    finally:
        if wav_path and os.path.exists(wav_path):
            os.remove(wav_path)

def extract_video(path):
    audio_tmp = None
    try:
        video = VideoFileClip(path)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            audio_tmp = tmp.name
        video.audio.write_audiofile(audio_tmp, logger=None)
        text = extract_audio(audio_tmp)
        return f"[Video]: {text}"
    except:
        return "[Video]: Không thể xử lý"
    finally:
        try:
            if 'video' in locals(): video.close()
            if audio_tmp and os.path.exists(audio_tmp): os.remove(audio_tmp)
        except:
            pass

# -------------------------
# MAIN DISPATCHER
# -------------------------
def extract_universal_content(path):
    ext = path.lower().split('.')[-1] if '.' in path else ''
    # Audio
    if ext in ['mp3', 'wav', 'ogg', 'flac', 'm4a', 'wma']:
        return extract_audio(path)
    # Video
    if ext in ['mp4', 'avi', 'mkv', 'mov', 'wmv', 'flv']:
        return extract_video(path)
    # Documents
    if ext == 'pdf': return extract_pdf(path)
    if ext in ['docx', 'doc']: return extract_docx(path)
    if ext in ['pptx', 'ppt']: return extract_pptx(path)
    if ext in ['xlsx', 'xls']: return extract_excel(path)
    if ext == 'csv': return extract_csv(path)
    # Structured / Code
    if ext == 'json': return extract_json(path)
    if ext in ['yaml', 'yml']: return extract_yaml(path)
    if ext in ['html', 'htm', 'xml']: return extract_html(path)
    if ext in ['sqlite', 'db']: return extract_sqlite(path)
    if ext == 'zip': return extract_zip(path)
    if ext in ['txt', 'md', 'py', 'js', 'ts', 'java', 'cpp', 'c', 'cs', 'php', 'css', 'log', 'ini', 'env']:
        return extract_txt(path)
    # Image fallback
    if USE_MAGIC:
        try:
            mime = magic.from_file(path, mime=True)
            if mime.startswith('image/'): return extract_image(path)
        except:
            pass
    return extract_from_binary(path)

# Tương thích cũ
extract_text = extract_universal_content
